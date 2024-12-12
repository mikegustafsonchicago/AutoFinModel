# -*- coding: utf-8 -*-
"""
Created on Mon Nov 18 13:00:59 2024

@author: mikeg

This module handles PowerPoint presentation generation for financial analysis.
It creates slides for revenue sources, costs, employees, expenses and capital expenditures.

Key components:
- Utility functions for styling and color handling
- Main presentation generation function
- Individual slide generation for each financial aspect
"""

#------------------------------------------------------------------------------
# Imports and Setup
#------------------------------------------------------------------------------
import os
import json
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.dml.color import MSO_THEME_COLOR
import sys
from pathlib import Path
sys.path.append(str(Path(__file__).parent.parent / 'excel_generation'))
from business_entity_code import BusinessEntity

#------------------------------------------------------------------------------
# Utility Functions
#------------------------------------------------------------------------------
def get_color(hex_code):
    """
    Convert hex color code to RGB tuple.
    
    Args:
        hex_code (str): Hex color code (e.g. '#FF0000')
    Returns:
        tuple: RGB values as (r,g,b)
    """
    hex_code = hex_code.lstrip('#')
    return int(hex_code[:2], 16), int(hex_code[2:4], 16), int(hex_code[4:], 16)

def load_json(filepath):
    """Load and parse JSON file."""
    with open(filepath, 'r') as file:
        return json.load(file)

def apply_table_styles(table, header_font_size=14):
    """
    Apply consistent styling to PowerPoint table headers and cells.
    Uses black and white table style.
    
    Args:
        table: PowerPoint table object
        header_font_size (int): Font size for header row
    """
    # Apply black and white table style
    table.style = 'Table Grid'
    
    # Style headers
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 0, 0)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(header_font_size)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
    
    # Style data cells with alternating shading
    for i in range(1, len(table.rows)):
        for cell in table.rows[i].cells:
            cell.fill.solid()
            # Alternate between white and light gray
            if i % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            else:
                cell.fill.fore_color.rgb = RGBColor(242, 242, 242)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)

#------------------------------------------------------------------------------
# Main Presentation Generation
#------------------------------------------------------------------------------
def generate_ppt():
    """
    Generate complete PowerPoint presentation with financial analysis.
    Creates slides for:
    - Title page
    - Revenue sources
    - Cost of sales
    - Employee overview
    - Operating expenses
    - Capital expenditures
    """
    current_directory = os.getcwd()
    file_name = "output_ppt.pptx"
    data_path = "temp_business_data"
    output_path = os.path.join(current_directory, "outputs", file_name)
    data_path = os.path.join(current_directory, data_path)
    
    # Load data using BusinessEntity
    try:
        business = BusinessEntity("financials")
        employees = {"employees": business.employees}
        revenue = {"revenue_sources": business.revenue_sources}
        cost_of_sales = {"cost_items": business.cost_of_sales_items}
        operating_expenses = {"expenses": business.operating_expenses}
        capex = {"expenses": business.capex_items}
    except Exception as e:
        print(f"Error loading business data: {str(e)}")
        raise

    # Create PowerPoint presentation
    presentation = Presentation()
    
    #--------------------------------------------------------------------------
    # Generate Individual Slides
    #--------------------------------------------------------------------------
    
    # Title slide
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "Business Plan Financial Analysis"
    subtitle.text = "Financial Projections and Analysis"

    # Revenue Sources slide
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Revenue Sources"
    
    rows = len(revenue["revenue_sources"]) + 1
    cols = 4
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(4)).table

    headers = ["Revenue Stream", "Price", "Monthly Transactions", "Source"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header

    for i, source in enumerate(revenue["revenue_sources"], 1):
        table.cell(i, 0).text = source["revenue_source_name"]
        table.cell(i, 1).text = f"${source['revenue_source_price']}"
        table.cell(i, 2).text = str(source["monthly_transactions"])
        table.cell(i, 3).text = source["price_source"]

    apply_table_styles(table)

    # Cost of Sales slide
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Cost of Sales"

    rows = len(cost_of_sales["cost_items"]) + 1
    cols = 4
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(4)).table

    headers = ["Cost Item", "Cost per Unit", "Monthly Units", "Source"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header

    for i, item in enumerate(cost_of_sales["cost_items"], 1):
        table.cell(i, 0).text = item["cost_item_name"]
        table.cell(i, 1).text = f"${item['cost_per_unit']}"
        table.cell(i, 2).text = str(item["monthly_transactions"])
        table.cell(i, 3).text = item["cost_source"]

    apply_table_styles(table)

    # Employees slide
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Firm Employees"

    rows = len(employees["employees"]) + 1
    cols = 4
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(4)).table

    headers = ["Role", "Number", "Wage", "Monthly Hours"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header

    for i, employee in enumerate(employees["employees"], 1):
        table.cell(i, 0).text = employee["role"]
        table.cell(i, 1).text = str(employee["number"])
        table.cell(i, 2).text = f"${employee['wage']}/{'hr' if employee['wage_type']=='hourly' else 'yr'}"
        table.cell(i, 3).text = str(employee["monthly_hours"])

    apply_table_styles(table)

    # Operating Expenses slide
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Operating Expenses"

    rows = len(operating_expenses["expenses"]) + 1
    cols = 3
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(3)).table

    headers = ["Expense Name", "Amount", "Frequency"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header

    for i, expense in enumerate(operating_expenses["expenses"], 1):
        table.cell(i, 0).text = expense["expense_name"]
        table.cell(i, 1).text = f"${expense['amount']}"
        table.cell(i, 2).text = expense["frequency"]

    apply_table_styles(table)

    # Capital Expenditures slide
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Capital Expenditures"

    rows = len(capex["expenses"]) + 1
    cols = 4
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(3)).table

    headers = ["Item", "Amount", "Purchase Year", "Depreciation Life"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header

    for i, expense in enumerate(capex["expenses"], 1):
        table.cell(i, 0).text = expense["expense_name"]
        table.cell(i, 1).text = f"${expense['amount']}"
        table.cell(i, 2).text = str(expense["purchase_year"])
        table.cell(i, 3).text = f"{expense['depreciation_life']} years"

    apply_table_styles(table)


    # Historical Financials slide
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Historical Financial Performance"

    # Create table for historical financials
    rows = len(business.historical_financials) + 1
    cols = 6
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(3)).table

    # Set headers
    headers = ["Year", "Revenue", "EBITDA", "EBIT", "Net Income", "EBITDA Margin"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header

    # Populate data
    for i, year_data in enumerate(business.historical_financials, 1):
        table.cell(i, 0).text = str(year_data["year"])
        table.cell(i, 1).text = f"${year_data['revenue']:,.0f}"
        table.cell(i, 2).text = f"${year_data['ebitda']:,.0f}"
        table.cell(i, 3).text = f"${year_data['ebit']:,.0f}"
        table.cell(i, 4).text = f"${year_data['net_income']:,.0f}"
        # Calculate EBITDA margin as percentage
        ebitda_margin = (year_data['ebitda'] / year_data['revenue']) * 100
        table.cell(i, 5).text = f"{ebitda_margin:.1f}%"

    apply_table_styles(table)

    # Save PowerPoint
    presentation.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    generate_ppt()