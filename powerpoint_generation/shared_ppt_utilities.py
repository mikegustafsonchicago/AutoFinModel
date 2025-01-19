import os
import json
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor, MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pathlib import Path


#------------------------------------------------------------------------------
# Utility Functions
#------------------------------------------------------------------------------
def get_color(hex_code):
    """
    Convert hex color code to RGB tuple.
    
    Args:
        hex_code (str): Hex color code (e.g. '#FF0000')
    Returns:
        tuple: RGB values as (r,g,b)
    """
    hex_code = hex_code.lstrip('#')
    return int(hex_code[:2], 16), int(hex_code[2:4], 16), int(hex_code[4:], 16)

def load_json(filepath):
    """Load and parse JSON file."""
    with open(filepath, 'r') as file:
        return json.load(file)
    
def load_text(filepath):
    """
    Load and read text file content.
    
    Args:
        filepath (str): Path to text file
    Returns:
        str: Content of text file
    """
    with open(filepath, 'r') as file:
        return file.read()


def apply_table_styles(table, include_sources=True, header_font_size=11):
    """
    Apply consistent styling to PowerPoint tables and handle source links properly.
    
    Args:
        table: PowerPoint table object
        include_sources (bool): Whether to process source objects into hyperlinks
        header_font_size (int): Font size for header row
    """
    # Apply basic table style with black borders
    table.style = 'Table Grid'
    
    # Apply black borders to all cells
    for row in table.rows:
        for cell in row.cells:
            for border in ['top', 'right', 'bottom', 'left']:
                border_obj = getattr(cell.border, border)
                border_obj.color.rgb = RGBColor(0, 0, 0)  # Black
                border_obj.width = Pt(1)
    
    # Style headers - white background with black text
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(header_font_size)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text
        
        # Add padding through text frame margins
        text_frame = cell.text_frame
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        
        if include_sources:
            _process_source_object(cell)
    
    # Style data cells - white background with black text
    for i in range(1, len(table.rows)):
        for cell in table.rows[i].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
            
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)
            paragraph.font.name = 'Calibri'
            paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text
            
            # Add consistent cell padding
            text_frame = cell.text_frame
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.05)
            text_frame.margin_bottom = Inches(0.05)
            
            if include_sources:
                _process_source_object(cell)
    
    # Set standard column widths if not already set
    for col in table.columns:
        if col.width < Inches(1):
            col.width = Inches(2)

def _process_source_object(cell):
    """Helper function to process source objects into hyperlinks"""
    if isinstance(cell.text, str) and cell.text.startswith("{'object_type': 'source_object'"):
        try:
            source_data = eval(cell.text)  # Convert string to dict
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.text = ""  # Clear existing text
            run = paragraph.add_run()
            run.text = source_data.get('display_value', 'Source')
            if source_data.get('url') and source_data['url'] != '-No Source-':
                run.hyperlink.address = source_data['url']
            run.font.color.rgb = RGBColor(0, 112, 192)  # Blue color for links
            run.font.underline = True
        except:
            # If there's any error parsing, leave as plain text
            paragraph.text = source_data.get('display_value', 'Source')

def generate_template_diagnostics(presentation, layout_index=0):
    """Generate diagnostic information about the template."""
    slide = presentation.slides.add_slide(presentation.slide_layouts[layout_index])  # Use Title and Content layout
    
    # Try to use title placeholder first, otherwise add title textbox
    title_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 1:  # 1 is title placeholder
            title_shape = shape
            break
            
    if title_shape is None:
        # No title placeholder found, create textbox
        title_shape = slide.shapes.add_textbox(
            left=Inches(1),
            top=Inches(0.5),
            width=Inches(8),
            height=Inches(0.5)
        )
        
    title_frame = title_shape.text_frame
    title_frame.text = "Template Diagnostics"
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.size = Pt(24)
    
    # Get template dimensions
    dims = analyze_template_dimensions(presentation)
    
    # Find body placeholder
    body_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 2:  # 2 is body placeholder
            body_shape = shape
            break
            
    # If no body placeholder found, add content directly to slide
    if body_shape is None:
        y_pos = Inches(1.5)
        for i, layout in enumerate(presentation.slide_layouts):
            text = slide.shapes.add_textbox(Inches(1), y_pos, Inches(8), Inches(0.5))
            p = text.text_frame.add_paragraph()
            p.text = f"Layout {i}: {layout.name}"
            p.font.bold = True
            y_pos += Inches(0.3)
            
            # List all placeholders in this layout
            for ph in layout.placeholders:
                text = slide.shapes.add_textbox(Inches(1.5), y_pos, Inches(8), Inches(0.3))
                p = text.text_frame.add_paragraph()
                p.text = (f"→ Placeholder {ph.placeholder_format.idx}: {ph.name} "
                         f"(type: {ph.placeholder_format.type})")
                if hasattr(ph, 'left'):
                    p.text += f" at ({ph.left/914400:.2f}\", {ph.top/914400:.2f}\")"
                y_pos += Inches(0.25)
            
            y_pos += Inches(0.2)  # Add space between layouts
    else:
        # Add content to body placeholder
        text_frame = body_shape.text_frame
        text_frame.clear()  # Clear any existing text
        
        for i, layout in enumerate(presentation.slide_layouts):
            p = text_frame.add_paragraph()
            p.text = f"Layout {i}: {layout.name}"
            p.font.bold = True
            
            # List all placeholders in this layout
            for ph in layout.placeholders:
                p = text_frame.add_paragraph()
                p.level = 1  # Indent one level
                p.text = (f"→ Placeholder {ph.placeholder_format.idx}: {ph.name} "
                         f"(type: {ph.placeholder_format.type})")
                if hasattr(ph, 'left'):
                    p.text += f" at ({ph.left/914400:.2f}\", {ph.top/914400:.2f}\")"
            
            text_frame.add_paragraph()  # Add blank line between layouts
    
    return dims


def get_gallery_images(gallery_path):
    """
    Get list of image filenames from the gallery directory.
    
    Args:
        gallery_path (str): Path to gallery directory
    Returns:
        list: List of image filenames in the gallery directory
    """
    try:
        # Convert to proper path object and resolve
        gallery_path = Path(gallery_path).resolve()
        logging.info(f"Looking for gallery at: {gallery_path}")
        
        if not gallery_path.exists():
            logging.error(f"Gallery path does not exist: {gallery_path}")
            return []
            
        # Get list of files in gallery directory
        image_files = []
        for filename in gallery_path.iterdir():
            # Check if file is an image by extension
            if filename.suffix.lower() in ('.png', '.jpg', '.jpeg', '.gif', '.bmp'):
                image_files.append(filename.name)
        
        logging.info(f"Found {len(image_files)} images in gallery")
        return sorted(image_files)
        
    except Exception as e:
        logging.error(f"Error getting gallery images: {str(e)}")
        return []  # Return empty list instead of raising to avoid presentation failure


def estimate_text_width(text, font_size, font_name='Calibri', bold=False):
    """
    Estimate the width of text in PowerPoint given font parameters.
    Returns width in points.
    
    Args:
        text (str): Text to measure
        font_size (int): Font size in points
        font_name (str): Font family name
        bold (bool): Whether text is bold
    
    Returns:
        float: Estimated width in points
    """
    # Average character widths as proportion of font size
    # These are approximate multipliers based on common fonts
    char_width_ratios = {
        'Calibri': {
            'normal': {
                'uppercase': 0.6,
                'lowercase': 0.5,
                'digits': 0.5,
                'spaces': 0.3
            },
            'bold': {
                'uppercase': 0.65,
                'lowercase': 0.55,
                'digits': 0.55,
                'spaces': 0.3
            }
        }
    }
    
    # Default to Calibri if font not in our ratio dictionary
    ratios = char_width_ratios.get(font_name, char_width_ratios['Calibri'])
    ratios = ratios['bold'] if bold else ratios['normal']
    
    width = 0
    for char in text:
        if char.isupper():
            width += ratios['uppercase']
        elif char.islower():
            width += ratios['lowercase']
        elif char.isdigit():
            width += ratios['digits']
        elif char.isspace():
            width += ratios['spaces']
        else:
            # For symbols and other characters, use average of upper and lowercase
            width += (ratios['uppercase'] + ratios['lowercase']) / 2
            
    return width * font_size


def analyze_template_dimensions(presentation):
    """
    Analyze template to find standard margins based on title and footer ribbon placement.
    Checks main master slide first, then layouts if needed.
    
    Args:
        presentation: PowerPoint presentation object
    Returns:
        dict: Contains key dimensions like title_height, footer_height, and usable_area
    """
    dimensions = {
        'slide_width': presentation.slide_width,
        'slide_height': presentation.slide_height,
        'title_height': None,
        'footer_height': None,
        'footer_top': None
    }
    
    # First check the main master slide
    main_master = presentation.slide_masters[0]  # The first slide master is the main one
    
    logging.info("Analyzing main master slide...")
    
    # Check placeholders on main master
    title_found = False
    footer_found = False
    
    for placeholder in main_master.placeholders:
        # Check for title placeholder
        if placeholder.placeholder_format.type == 1:  # 1 = Title
            dimensions['title_height'] = int(placeholder.height)
            dimensions['title_top'] = placeholder.top
            title_found = True
            logging.info(f"Found title in master: height={placeholder.height}, top={placeholder.top}")
        
        # Look for footer placeholder
        if (placeholder.top + placeholder.height) > (presentation.slide_height - Inches(1)):
            footer_top = placeholder.top
            footer_height = int(placeholder.height)
            footer_found = True
            logging.info(f"Found footer in master: height={placeholder.height}, top={placeholder.top}")
    
    # Check shapes in main master
    bottom_area = {
        'top': None,
        'height': None
    }
    
    bottom_threshold = presentation.slide_height - Inches(1.5)
    for shape in main_master.shapes:
        # Check if shape is in bottom area
        if shape.top > bottom_threshold:
            shape_bottom = shape.top + shape.height
            logging.info(f"Found bottom shape in master: top={shape.top}, height={shape.height}, type={type(shape)}")
            
            # Initialize or update bottom area bounds
            if bottom_area['top'] is None or shape.top < bottom_area['top']:
                bottom_area['top'] = shape.top
            
            if bottom_area['height'] is None or shape_bottom > (bottom_area['top'] + bottom_area['height']):
                bottom_area['height'] = shape_bottom - bottom_area['top']
    
    # If we didn't find what we need in the master, check layouts
    if not title_found or (not footer_found and bottom_area['top'] is None):
        logging.info("Required elements not found in master, checking layouts...")
        # [Previous layout checking code here]
        for layout in presentation.slide_layouts:
            for placeholder in layout.placeholders:
                if not title_found and placeholder.placeholder_format.type == 1:
                    dimensions['title_height'] = int(placeholder.height)
                    dimensions['title_top'] = placeholder.top
                    title_found = True
                
                if not footer_found and (placeholder.top + placeholder.height) > (presentation.slide_height - Inches(1)):
                    footer_top = placeholder.top
                    footer_height = int(placeholder.height)
                    footer_found = True
            
            # Check shapes in layout if we still need bottom area
            if bottom_area['top'] is None:
                for shape in layout.shapes:
                    if shape.top > bottom_threshold:
                        shape_bottom = shape.top + shape.height
                        if bottom_area['top'] is None or shape.top < bottom_area['top']:
                            bottom_area['top'] = shape.top
                        if bottom_area['height'] is None or shape_bottom > (bottom_area['top'] + bottom_area['height']):
                            bottom_area['height'] = shape_bottom - bottom_area['top']
    
    # Determine which measurement to use for footer area
    if bottom_area['top'] is not None and bottom_area['height'] is not None:
        logging.info(f"Using bottom area: top={bottom_area['top']}, height={bottom_area['height']}")
        # Use bottom area if it's higher up than footer placeholder
        if footer_found and bottom_area['top'] < footer_top:
            dimensions['footer_top'] = bottom_area['top']
            dimensions['footer_height'] = bottom_area['height']
        else:
            dimensions['footer_top'] = footer_top
            dimensions['footer_height'] = footer_height
    elif footer_found:
        logging.info("Using footer placeholder measurements")
        # Fall back to footer placeholder if no bottom shapes found
        dimensions['footer_top'] = footer_top
        dimensions['footer_height'] = footer_height
    else:
        raise ValueError("Could not find required template elements (title and footer/bottom area)")
    
    # Calculate usable content area
    if dimensions['title_height'] and dimensions['footer_height']:
        dimensions['content_top'] = int(dimensions['title_top'] + dimensions['title_height'])   
        dimensions['content_height'] = int(dimensions['footer_top'] - dimensions['content_top'])
        dimensions['content_width'] = presentation.slide_width
        
        # Add some helpful derived measurements
        dimensions['left_column_width'] = (dimensions['content_width'] / 2) - Inches(0.125)
        dimensions['right_column_width'] = dimensions['left_column_width']
        dimensions['column_spacing'] = Inches(0.25)
        
        
        return dimensions
    else:
        raise ValueError("Could not find required template elements (title and footer/bottom area)")



def add_diagnostics_slide(presentation):
    """
    Add a slide showing template dimensions and theme colors.
    
    Args:
        presentation: PowerPoint presentation object
    """
    # Create a slide to show dimensions and colors
    slide = presentation.slides.add_slide(presentation.slide_layouts[10])
    title = slide.shapes.title
    title.text = "Template Dimensions & Colors"
    
    # Get template dimensions
    dims = analyze_template_dimensions(presentation)
    
    # Add dimension visualization shapes
    # Title area
    title_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
        Inches(1), dims['title_top'], Inches(2), dims['title_height'])
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = RGBColor(200, 200, 200)
    title_box.line.color.rgb = RGBColor(100, 100, 100)
    
    # Add text to title box
    title_text = title_box.text_frame
    title_text.text = f"Title Area\nHeight: {dims['title_height']/914400:.2f}\"\nTop: {dims['title_top']/914400:.2f}\" (title_top)"
    title_text.paragraphs[0].font.size = Pt(8)
    title_text.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    title_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Content area
    content_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(1), dims['content_top'], Inches(2), dims['content_height'])
    content_box.fill.solid() 
    content_box.fill.fore_color.rgb = RGBColor(230, 230, 230)
    content_box.line.color.rgb = RGBColor(100, 100, 100)
    
    # Add text to content box
    content_text = content_box.text_frame
    content_text.text = f"Content Area\nHeight: {dims['content_height']/914400:.2f}\"\nTop: {dims['content_top']/914400:.2f}\" (content_top)"
    content_text.paragraphs[0].font.size = Pt(8)
    content_text.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    content_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Footer area
    footer_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(1), dims['footer_top'], Inches(2), dims['footer_height'])
    footer_box.fill.solid()
    footer_box.fill.fore_color.rgb = RGBColor(200, 200, 200)
    footer_box.line.color.rgb = RGBColor(100, 100, 100)
    
    # Add text to footer box
    footer_text = footer_box.text_frame
    footer_text.text = f"Footer Area\nHeight: {dims['footer_height']/914400:.2f}\"\nTop: {dims['footer_top']/914400:.2f}\" (footer_top)"
    footer_text.paragraphs[0].font.size = Pt(10)
    footer_text.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    footer_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    footer_text.paragraphs[1].font.size = Pt(10)
    footer_text.paragraphs[1].font.color.rgb = RGBColor(0, 0, 0)
    footer_text.paragraphs[1].alignment = PP_ALIGN.CENTER
    
    # Add dimension labels
    labels = {
        'Title Height': dims['title_height'],
        'Content Top': dims['content_top'],
        'Content Height': dims['content_height'],
        'Footer Height': dims['footer_height'],
        'Footer Top': dims['footer_top'],
        'Slide Width': presentation.slide_width,
        'Slide Height': presentation.slide_height
    }
    
    # Add labels with measurements
    for i, (label, value) in enumerate(labels.items()):
        text = slide.shapes.add_textbox(Inches(4), Inches(1.5 + i*0.4), Inches(4), Inches(0.3))
        text.text_frame.text = f"{label}: {value/914400:.2f} inches"  # Convert EMU to inches
        
    # Add theme colors in compact form
    theme_colors = {
        'ACCENT_1': MSO_THEME_COLOR.ACCENT_1,
        'ACCENT_2': MSO_THEME_COLOR.ACCENT_2,
        'DARK_1': MSO_THEME_COLOR.DARK_1,
        'LIGHT_1': MSO_THEME_COLOR.LIGHT_1,
    }
    
    for i, (name, color) in enumerate(theme_colors.items()):
        shape = slide.shapes.add_textbox(Inches(4), Inches(5 + i*0.3), Inches(4), Inches(0.3))
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.text = f"Color: {name}"
        paragraph.font.theme_color = color
