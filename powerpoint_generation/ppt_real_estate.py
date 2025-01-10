# -*- coding: utf-8 -*-
"""
Created on Mon Nov 18 13:00:59 2024

@author: mikeg

This module handles PowerPoint presentation generation for real estate property analysis.
It creates slides showing property details, zoning information, and financial metrics.

Key components:
- Utility functions for styling and color handling 
- Main presentation generation function
- Individual slide generation for property aspects
"""

#------------------------------------------------------------------------------
# Imports and Setup
#------------------------------------------------------------------------------
import sys
import os
sys.path.append(os.path.dirname(os.path.dirname(__file__)))
import json
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.dml.color import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, UnidentifiedImageError
from file_manager import (
    get_project_outputs_path, 
    get_project_gallery_path, 
    get_project_data_path
)
from .shared_ppt_utilities import (
    load_json,
    get_gallery_images,
    estimate_text_width,
    analyze_template_dimensions,
    apply_table_styles,
    generate_template_diagnostics,
    add_diagnostics_slide
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(levelname)s - %(message)s'
)

#------------------------------------------------------------------------------
# Main Presentation Generation
#------------------------------------------------------------------------------
def generate_real_estate_ppt():

    class PresentationData:
        def __init__(self):
            self.title_image_name = "southcarolina-header.jpg"
            self.file_name = "property_presentation.pptx"
            self.template_name = os.path.join(os.path.dirname(__file__), "hucks_template.pptx")
            self.output_path = get_project_outputs_path()
            self.data_path = get_project_data_path()
            self.gallery_path = get_project_gallery_path()
            self.gallery_images = get_gallery_images(self.gallery_path)
            
            # Load property data and firm data
            try:
                self.property_fundamentals = load_json(os.path.join(self.data_path, "property_fundamentals.json"))
                self.property_zoning = load_json(os.path.join(self.data_path, "property_zoning.json"))
                self.property_financials = load_json(os.path.join(self.data_path, "property_financials.json"))
                self.firm_data = load_json(os.path.join(os.path.dirname(__file__), "firm_data.json"))
            except Exception as e:
                logging.error(f"Error loading property or firm data: {str(e)}")
                raise
            
    ppt_data = PresentationData()
    gallery_running_index = 0

    # Create PowerPoint presentation
    presentation = Presentation(ppt_data.template_name)
    
    # Create template dimensions class to pass to slide functions
    class TemplateDimensions:
        def __init__(self, presentation):
            self.slide_width = presentation.slide_width
            self.slide_height = presentation.slide_height
            
            # Analyze template dimensions
            template_dims = analyze_template_dimensions(presentation)
            self.title_height = template_dims['title_height']
            self.footer_height = template_dims['footer_height'] 
            self.content_top = template_dims['content_top']
            self.content_height = template_dims['content_height']
            self.content_width = template_dims['content_width']
            
    # Create single instance to use throughout
    template_dims = TemplateDimensions(presentation)

    logging.info(f"Template Dimensions:\n Title Height: {template_dims.title_height}, Footer Height: {template_dims.footer_height}")
    


    #--- 1 ---Title Slide---
    try:
        add_title_slide(presentation, ppt_data, template_dims)
    
    except Exception as e:
        logging.error(f"Error creating title slide: {str(e)}")

    #--- 2 --- Property Fundamentals Slide---
    try:
        add_property_fundamentals_slide(presentation, ppt_data, template_dims)
    except Exception as e:
        logging.error(f"Error creating property fundamentals slide: {str(e)}")

    
    #--- 3 --- Zoning & Financial Details Slide---
    try:
        add_zoning_financial_details_slide(presentation, ppt_data, template_dims)
    except Exception as e:
        logging.error(f"Error creating zoning & financial details slide: {str(e)}")

    
    #--- 4 through end-2 --- Photo Gallery Slides---
    try:
        add_photo_gallery_slides(presentation, ppt_data, template_dims)
    except Exception as e:
        logging.error(f"Error creating photo gallery slides: {str(e)}")

    
    #--- Last-1 --- About Us Slide---
    try:
        add_about_us_slide(presentation, ppt_data, template_dims)
    except Exception as e:
        logging.error(f"Error creating about us slide: {str(e)}")

    
    #--- Last --- Legal Disclaimer Slide---
    try:
        add_legal_disclaimer_slide(presentation, ppt_data, template_dims)
    except Exception as e:
        logging.error(f"Error creating legal disclaimer slide: {str(e)}")


    #---Last +1 --- Diagnostics Slide---
    try:
        add_diagnostics_slide(presentation, ppt_data, template_dims)
    except Exception as e:
        logging.error(f"Error creating diagnostics slide: {str(e)}")


    # Generate template diagnostics text file
    try:
        generate_template_diagnostics(ppt_data.template_name, ppt_data.output_path)
    except Exception as e:
        logging.error(f"Error creating diagnostics text: {str(e)}")

    # Save PowerPoint
    output_file = os.path.join(ppt_data.output_path, ppt_data.file_name)
    presentation.save(output_file)
    logging.info(f"Presentation saved to {ppt_data.output_path}")
    return ppt_data.file_name

if __name__ == "__main__":
    generate_real_estate_ppt()




#--------------------------------------------------------------------------
# SLIDE 1: Title Slide
#--------------------------------------------------------------------------

def add_title_slide(presentation, ppt_data, presentation_dims):       
    # Use blank layout instead of title layout
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # Usually layout 6 is blank

    # Add background image first
    background_img_path = os.path.join(ppt_data.gallery_path, ppt_data.title_image_name)
    if not os.path.exists(background_img_path):
        logging.warning(f"Title slide background image not found at {background_img_path}. Using default background.")
        background_img_path = os.path.join(os.path.dirname(__file__), "default_background.jpg")

    left = Inches(0)
    top = Inches(0)
    background = title_slide.shapes.add_picture(
        background_img_path, 
        left, 
        top, 
        presentation_dims.slide_width, 
        presentation_dims.slide_height
    )

    # Add text boxes manually instead of using placeholders
    # Calculate text width first to center the box
    title_text = f"Investment Overview: {ppt_data.property_fundamentals['property_details'][0]['address']}"
    text_width = estimate_text_width(title_text, Pt(44))
    text_width = int(max(text_width, Inches(4)))  # Minimum 4 inches

    # Calculate left position to center the text box
    left = (presentation_dims.slide_width - text_width) / 2
    title_box = title_slide.shapes.add_textbox(left, Inches(2), text_width, Inches(1.5))

    title_frame = title_box.text_frame
    title_para = title_frame.add_paragraph()
    title_para.text = title_text
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER  # Center align the text within the box

    # Add white outline
    title_box.line.color.rgb = RGBColor(255, 255, 255)
    title_box.line.width = Pt(2)

    line = title_box.line
    line.color.rgb = RGBColor(255, 255, 255)
    line.width = Pt(2)

    subtitle_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(8), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.add_paragraph()
    subtitle_para.text = f"{ppt_data.property_fundamentals['property_details'][0]['municipality']}\n\n{ppt_data.firm_data['firm_details']['name']}\n{ppt_data.firm_data['firm_details']['contacts'][0]['email']}"
    subtitle_para.font.size = Pt(14)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)

#--------------------------------------------------------------------------
# SLIDE 2: Property Fundamentals
#--------------------------------------------------------------------------
def add_property_fundamentals_slide(presentation, ppt_data, presentation_dims):   
    try:
        slide = presentation.slides.add_slide(presentation.slide_layouts[10])
        title = slide.shapes.title
        title.text = "Property Highlights"
        
        fund_details = ppt_data.property_fundamentals['property_details'][0]
        
        # Get fields from property_fundamentals_structure.json
        display_fields = []
        properties = ppt_data.property_fundamentals_structure['property_fundamentals']['structure']['property_details']['items']['properties']
        
        # Add fields in order, skipping source/notes fields
        for i in range(1, len(properties) + 1):
            for field, details in properties.items():
                if details.get('order') == i and not field.endswith(('_source', '_notes')):
                    display_fields.append(field)
        
        rows = len(display_fields) + 1  # +1 for header
        cols = 2
        
        # Calculate table and image widths based on split ratio (e.g. 40% table, 60% image)
        table_width_ratio = 0.4  # Adjust this value to change the split
        image_width_ratio = 1 - table_width_ratio
        
        table_width = int(presentation_dims.slide_width * table_width_ratio - Inches(0.25))  # Subtract margin
        
        # Create table with new width
        table = slide.shapes.add_table(rows, cols, Inches(0.125), Inches(1.5),
                                    table_width, Inches(4.00)).table

        # Merge header cells and set text
        table.cell(0, 0).merge(table.cell(0, 1))
        table.cell(0, 0).text = "Overview"

        row = 1
        for field in display_fields:
            # Set white background for each row
            for col in range(cols):
                cell = table.cell(row, col)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            # Get value, skipping source/notes fields
            value = fund_details[field]
            if isinstance(value, dict) and value.get('object_type') == 'source_object':
                continue
                
            table.cell(row, 0).text = field.replace('_', ' ').title()
            table.cell(row, 1).text = str(value)
            row += 1

        apply_table_styles(table)

        # Add gallery image to right side of slide
        if ppt_data.gallery_images[gallery_running_index]:
            selected_image = os.path.join(ppt_data.gallery_path, ppt_data.gallery_images[gallery_running_index])
            # Get image dimensions before adding to slide
            with Image.open(selected_image) as img:
                img_width, img_height = img.size
                # Calculate aspect ratio
                aspect_ratio = img_width / img_height
            
            # Calculate image position and size based on ratio
            left_x = int(presentation_dims.slide_width * table_width_ratio + Inches(0.125))  # Add margin
            image_width = int(presentation_dims.slide_width * image_width_ratio - Inches(0.25))  # Subtract margin
            image_height = image_width / aspect_ratio
            
            # Center image vertically
            top_y = presentation_dims.title_height + ((presentation_dims.slide_height - presentation_dims.title_height - presentation_dims.footer_height - image_height) / 2)
        
            # Add image
            image = slide.shapes.add_picture(selected_image, left_x, top_y,
                                        image_width, image_height)
            gallery_running_index += 1
    except Exception as e:
        logging.error(f"Failed to add property fundamentals slide: {str(e)}")
        # Continue without this slide rather than failing presentation generation



#--------------------------------------------------------------------------
# SLIDE 3: Zoning & Financial Details
#--------------------------------------------------------------------------
def add_zoning_financial_details_slide(presentation, ppt_data, presentation_dims):
    try:
        slide = presentation.slides.add_slide(presentation.slide_layouts[10])
        title = slide.shapes.title
        title.text = "Zoning & Financial Details"
            
        # Left table - Zoning details
        zoning_details = ppt_data.property_zoning['property_details']
        rows = len(zoning_details) + 1
        cols = 2

        width = int((presentation_dims.slide_width/2)*.8)  
        left_x = int(((presentation_dims.slide_width/2) - width)/2)  
        top_y = Inches(1.5)
        height = Inches(4.00)
        left_table = slide.shapes.add_table(rows, cols, left_x, top_y, width, height).table
        logging.debug("-" * 80)
        logging.debug(f"Left Table Dimensions:")
        logging.debug(f"Left: {left_x}")
        logging.debug(f"Top: {top_y}")
        logging.debug(f"Width: {width}")
        logging.debug(f"Height: {height}")
        logging.debug(f"Slide width: {presentation_dims.slide_width}")
        logging.debug("-" * 80)

        # Merge header cells and set text
        left_table.cell(0, 0).merge(left_table.cell(0, 1))
        left_table.cell(0, 0).text = "Zoning Overview"

        row = 1
        for key, value in zoning_details.items():
            left_table.cell(row, 0).text = key.replace('_', ' ').title()
            left_table.cell(row, 1).text = str(value)
            row += 1

        apply_table_styles(left_table)

        # Right table - Financial details
        financial_details = ppt_data.property_financials['property_details']
        rows = len(financial_details) + 1
        cols = 2
        width = int((presentation_dims.slide_width/2)*.8) 
        left_x = int(((presentation_dims.slide_width/2) - width)/2+presentation_dims.slide_width/2)
        right_table = slide.shapes.add_table(rows, cols, left_x, top_y, width, height).table
        logging.debug("-" * 80)
        logging.debug(f"Right Table Dimensions:")
        logging.debug(f"Left: {left_x}")
        logging.debug(f"Top: {top_y}")
        logging.debug(f"Width: {width}")
        logging.debug(f"Height: {height}")
        logging.debug(f"Slide width: {presentation_dims.slide_width}")
        logging.debug("-" * 80)

        # Merge header cells and set text
        right_table.cell(0, 0).merge(right_table.cell(0, 1))
        right_table.cell(0, 0).text = "Financial Overview"

        row = 1
        for key, value in financial_details.items():
            right_table.cell(row, 0).text = key.replace('_', ' ').title()
            right_table.cell(row, 1).text = str(value)
            row += 1

        apply_table_styles(right_table)
    except Exception as e:
        logging.error(f"Failed to add property fundamentals slide: {str(e)}")
        # Continue without this slide rather than failing presentation generation 


#--------------------------------------------------------------------------
    # Photo Gallery Slides
    #--------------------------------------------------------------------------
    # Get list of images from gallery folder
def add_photo_gallery_slides(presentation, ppt_data, presentation_dims):
    
    try:
        if os.path.exists(ppt_data.gallery_path):
            image_files = [f for f in os.listdir(ppt_data.gallery_path) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]

            
            for image_file in image_files:
                if image_file == ppt_data.title_image_name:
                    continue
                # Add new slide for each image
                slide = presentation.slides.add_slide(presentation.slide_layouts[10])  # Blank layout
                
                # Add title
                title = slide.shapes.title
                title.text = "Property View"
                
                # Calculate dimensions to maintain aspect ratio while fitting in slide
                img_path = os.path.join(ppt_data.gallery_path, image_file)
                
                # Center the image horizontally and vertically

                
                # Get original image dimensions
                with Image.open(img_path) as img:
                    img_width, img_height = img.size
                
                # Calculate aspect ratio
                image_aspect_ratio = img_width / img_height
                full_area_ratio = presentation_dims.slide_width / (presentation_dims.slide_height - presentation_dims.title_height)

                logging.debug(f" Image Name: {image_file}")
                logging.debug(f" Image Aspect Ratio: {round(image_aspect_ratio, 2)}")
                logging.debug(f" Full Aspect Ratio: {round(full_area_ratio, 2)}")
                logging.debug(f" Delta: {round(abs(image_aspect_ratio - full_area_ratio), 2)}")
                logging.debug(f" SLIDE_HEIGHT: {presentation_dims.slide_height}, TITLE_HEIGHT: {presentation_dims.title_height}, FOOTER_HEIGHT: {presentation_dims.footer_height}")
                logging.debug("\n")
                
                # Set max dimensions while preserving aspect ratio
                max_width = presentation_dims.slide_width
                max_height = presentation_dims.slide_height - presentation_dims.title_height
                
                # First, check if image aspect ratio is suitable for the full slide area below title
                # This area includes the footer ribbon, which is fine if we want to cover it
                ratio_difference = abs(image_aspect_ratio - full_area_ratio)
                
                if ratio_difference < 0.45:
                    # Image ratio is close enough to full area ratio (within 40%)
                    # We'll stretch/compress slightly to fill the full area
                    # This is acceptable since the difference is small
                    width = max_width
                    height = max_height
                else:
                    # Image ratio differs too much from full area ratio
                    # Now we'll fit it in the content area between title and footer
                    content_area_ratio = presentation_dims.content_width / (presentation_dims.content_height)
                    
                    if image_aspect_ratio > content_area_ratio:
                        # Image is wider than content area
                        # Use full width and calculate proportional height
                        width = max_width
                        height = round(width / image_aspect_ratio)
                    else:
                        # Image is taller than content area
                        # Use available height (excluding footer) and calculate proportional width
                        height = presentation_dims.content_height
                        width = round(height * image_aspect_ratio)
                # Calculate centered position
                left = int((presentation_dims.slide_width - width) / 2)
                top_y = presentation_dims.content_top
        slide.shapes.add_picture(img_path, left, top_y, width, height)
    
    except Exception as e:
        logging.error(f"Failed to add property fundamentals slide: {str(e)}")
        # Continue without this slide rather than failing presentation generation
            


#--------------------------------------------------------------------------
# SLIDE Last-1: About Us
#--------------------------------------------------------------------------
def add_about_us_slide(presentation, ppt_data, presentation_dims):
    
    try:
        slide = presentation.slides.add_slide(presentation.slide_layouts[10])
        title = slide.shapes.title
        title.text = "About Us"

        # Add firm details from firm_data.json
        firm_details = ppt_data.firm_data['firm_details']
        
        # Add text box for firm details
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        # Add company description
        p = text_frame.add_paragraph()
        p.text = firm_details['about']
        p.font.size = Pt(11)
        p.font.name = 'Calibri'
        p.space_after = Pt(12)
        
        # Add contact information
        p = text_frame.add_paragraph()
        p.text = f"Address: {firm_details['address']}"
        p.font.size = Pt(10)
        p.font.name = 'Calibri'
        
        p = text_frame.add_paragraph()
        p.text = f"Website: {firm_details['website']}"
        p.font.size = Pt(10)
        p.font.name = 'Calibri'
        p.space_after = Pt(12)
        
        # Add team contacts
        for contact in firm_details['contacts']:
            p = text_frame.add_paragraph()
            p.text = f"{contact['name']} - {contact['title']}"
            p.font.size = Pt(10)
            p.font.name = 'Calibri'
            
            p = text_frame.add_paragraph()
            p.text = f"Phone: {contact['phone']} | Email: {contact['email']}"
            p.font.size = Pt(10)
            p.font.name = 'Calibri'
            p.space_after = Pt(8)
    
    except Exception as e:
        logging.error(f"Failed to add property fundamentals slide: {str(e)}")
        # Continue without this slide rather than failing presentation generation

#--------------------------------------------------------------------------
# SLIDE Last: Legal Disclaimer
#--------------------------------------------------------------------------
def add_legal_disclaimer_slide(presentation, ppt_data, presentation_dims):
    try:
        slide = presentation.slides.add_slide(presentation.slide_layouts[11])
        '''
        title = slide.shapes.title
        title.text = "Legal Disclaimer"

        # Load about us text
        about_us_path = os.path.join(os.path.dirname(__file__), "disclaimer.txt")
        with open(about_us_path, 'r') as f:
            about_us_text = f.read()

        # Add text box for about us content
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        # Format text by splitting into paragraphs and applying consistent formatting
        text_frame.clear()  # Clear existing text
        paragraphs = about_us_text.split('\n\n')  # Split on double newlines
        
        for text in paragraphs:
            if text.strip():  # Only process non-empty paragraphs
                p = text_frame.add_paragraph()
                p.text = text.strip()
                p.font.size = Pt(10)
                p.font.name = 'Calibri'
                p.space_after = Pt(12)  # Add spacing between paragraphs
    '''
    except Exception as e:
        logging.error(f"Failed to add property fundamentals slide: {str(e)}")
        # Continue without this slide rather than failing presentation generation