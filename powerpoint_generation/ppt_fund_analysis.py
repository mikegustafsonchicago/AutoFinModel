# -*- coding: utf-8 -*-
"""
Created on Mon Nov 18 13:00:59 2024

@author: mikeg

This module handles PowerPoint presentation generation for financial analysis.
It creates slides for revenue sources, costs, employees, expenses and capital expenditures.

Key components:
- Utility functions for styling and color handling
- Main presentation generation function
- Individual slide generation for each financial aspect
"""

#------------------------------------------------------------------------------
# Imports and Setup
#------------------------------------------------------------------------------
import os
import json
import yaml
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from datetime import datetime as dt
from io import BytesIO
import sys
from pathlib import Path
sys.path.append(str(Path(__file__).parent.parent / 'excel_generation'))
# Local imports
from business_entity_code import BusinessEntity
from file_manager import (
    get_project_outputs_path, 
    get_project_gallery_path, 
    get_project_data_path
)
from .shared_ppt_utilities import (
    load_json,
    get_gallery_images,
    estimate_text_width,
    analyze_template_dimensions,
    apply_table_styles,
    generate_template_diagnostics,
    add_diagnostics_slide
)
#------------------------------------------------------------------------------
# Utility Functions
#------------------------------------------------------------------------------
def get_color(hex_code):
    """
    Convert hex color code to RGB tuple.
    
    Args:
        hex_code (str): Hex color code (e.g. '#FF0000')
    Returns:
        tuple: RGB values as (r,g,b)
    """
    hex_code = hex_code.lstrip('#')
    return int(hex_code[:2], 16), int(hex_code[2:4], 16), int(hex_code[4:], 16)

def load_json(filepath):
    """Load and parse JSON file."""
    with open(filepath, 'r') as file:
        return json.load(file)
    

#------------------------------------------------------------------------------
# Main Presentation Generation
#------------------------------------------------------------------------------
def generate_fund_analysis_ppt(debug=False):

    class InvestmentFirmData:
        def __init__(self):
            # Basic firm info
            self.firm_name = "3 Pillars Capital"
            self.firm_address = "123 Main St, New York, NY 10001"
            self.firm_phone = "(212) 555-1234"
            self.firm_email = "info@3pillars.com"
            self.firm_website = "www.3pillars.com"
            
            # Investment details
            self.commitment_size = 5000000
            self.target_fund_size = 10000000
            self.investment_strategy = 'middle market distressed debt and mezzanine fund'
            self.previous_exposure = 'None'
            
            # Team info
            self.team_members = [
                {"name": "Amanda Lin", "title": "Managing Partner", "email": "amanda@3pillars.com"},
                {"name": "Mike Gustafson", "title": "Senior Associate", "email": "mike@3pillars.com"},
                {"name": "Jane Doe", "title": "Analyst", "email": "jane@3pillars.com"}
            ]

            self.primary_contact = "Mike Gustafson (212) 555-1234"
            
            # Recommendation details
            self.recommendation_amount = '$7.5MM'
            self.submission_date = dt.now().strftime('%B %d, %Y')
            self.next_closing = 'Early week of December 16, 2002'
            self.turnaround_requested = '48-hour'
            


    class PresentationData:
        def __init__(self):
            self.file_name = "fund_analysis.pptx"
            self.template_name = os.path.join(os.path.dirname(__file__), "Investment_Overview_Template.pptx")
            self.output_path = get_project_outputs_path()
            self.data_path = get_project_data_path()
            self.gallery_path = get_project_gallery_path()
            logging.debug(f"Gallery path: {self.gallery_path}")
            self.gallery_images = get_gallery_images(self.gallery_path)

        """Generate a PowerPoint presentation for fund analysis"""
        # Track debug info
    debug_info = {
        'errors': [],
        'slide_info': [],
        'data_loaded': {},
        'timing': {}
    }
    start_time = dt.now()

    investment_firm = InvestmentFirmData()
    ppt_data = PresentationData()
    
    # Create presentation from template
    try:
        prs = Presentation(ppt_data.template_name)
        # Log available layouts to help debug
        logging.info(f"Available slide layouts: {len(prs.slide_layouts)}")

    except Exception as e:
        logging.error(f"Error loading template: {str(e)}")
        raise

    # Now use correct layout indices
    LAYOUTS = {
        'title': 0,        # Usually the title slide layout
        'four_block': 1,      # Section header layout
        'table': 2,      # Table layout
        'table_content': 3,      # Table layout with content  
        'content': 4,      # Content layout with title and body
        'blank': 5,        # Blank slide layout
    }

    # Load data using BusinessEntity
    try:
        business = BusinessEntity("fund_analysis")
        fundamentals = business.fundamentals[0]  # Get first entry
        investment_team = business.investment_team
        seed_terms = business.seed_terms
        fees_key_terms = business.fees_key_terms
        deal_history = business.deal_history
        service_providers = business.service_providers
        
        debug_info['data_loaded'] = {
            'fundamentals': bool(fundamentals),
            'investment_team': len(investment_team),
            'seed_terms': bool(seed_terms),
            'fees_key_terms': bool(fees_key_terms),
            'deal_history': len(deal_history),
            'service_providers': len(service_providers)
        }
    except Exception as e:
        error_msg = f"Error loading fund data: {str(e)}"
        logging.error(error_msg)
        debug_info['errors'].append(error_msg)
        raise
    #------------------------------------------------------------------------------
    # Title Slide
    #------------------------------------------------------------------------------
    try:
        slide_start = dt.now()
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['title']])
        
        # Set the title to firm name
        title = slide.shapes.title
        title.text = investment_firm.firm_name
        
        # Create a dictionary of all placeholders by index
        placeholders = {shape.placeholder_format.idx: shape for shape in slide.placeholders}
        
        # Add firm name and commitment info (placeholder 14)
        if 14 in placeholders:
            commitment_frame = placeholders[14].text_frame
            commitment_frame.clear()
            
            # Firm name and commitment
            firm_para = commitment_frame.add_paragraph()
            firm_para.text = fundamentals.get('firm_name', investment_firm.firm_name)
            firm_para.font.size = Pt(24)
            firm_para.font.bold = True
            firm_para.space_after = Pt(12)
            
            commitment_para = commitment_frame.add_paragraph()
            commitment_para.text = f"Blind Pool Commitment of ${investment_firm.commitment_size:,} in a Private Equity Fund"
            commitment_para.font.size = Pt(16)
            
            # Add previous exposure here
            if investment_firm.previous_exposure:
                exposure_para = commitment_frame.add_paragraph()
                exposure_para.text = f"({investment_firm.previous_exposure})"
                exposure_para.font.size = Pt(12)
                exposure_para.font.italic = True
        
        # Add target fund size and strategy (placeholder 15)
        if 15 in placeholders:
            target_frame = placeholders[15].text_frame
            target_frame.clear()
            target_para = target_frame.add_paragraph()
            target_para.text = f"Sponsor is targeting ${investment_firm.target_fund_size:,} in commitments for this"
            target_para.font.size = Pt(14)
            target_para.font.italic = True
            
            strategy_para = target_frame.add_paragraph()
            strategy_para.text = investment_firm.investment_strategy
            strategy_para.font.size = Pt(14)
            strategy_para.font.italic = True
        
        # Create recommendation table instead of using placeholder
        left = Inches(7.07)  # Same as original placeholder
        top = Inches(2.99)   # Same as original placeholder
        width = Inches(4.5)  # Adjust as needed
        height = Inches(4)   # Adjust as needed
        
        table = slide.shapes.add_table(
            rows=6 + len(investment_firm.team_members),
            cols=2,
            left=left,
            top=top,
            width=width,
            height=height
        ).table
        
        # Row 1 - CF Recommendation
        table.cell(0, 0).text = "CF Recommendation:"
        table.cell(0, 0).text_frame.paragraphs[0].font.bold = True
        table.cell(0, 1).text = f"{investment_firm.recommendation_amount} LP Investment"
        
        # Row 2 - Disclaimer (merged cells)
        table.rows[1].cells[0].merge(table.rows[1].cells[1])
        cell = table.rows[1].cells[0]
        cell.text = "(Subject to Legal, Background and Tax Due Diligence)"
        cell.text_frame.paragraphs[0].font.italic = True
        
        # Row 3 - Team Header (merged cells)
        table.rows[2].cells[0].merge(table.rows[2].cells[1])
        cell = table.rows[2].cells[0]
        cell.text = "Equity Deal Team:"
        
        # Add team members
        current_row = 3
        for member in investment_firm.team_members:
            table.rows[current_row].cells[0].merge(table.rows[current_row].cells[1])
            cell = table.rows[current_row].cells[0]
            cell.text = f"{member['name']} - {member['title']} - {member['email']}"
            current_row += 1
        
        # Add dates (merged cells)
        table.rows[current_row].cells[0].merge(table.rows[current_row].cells[1])
        cell = table.rows[current_row].cells[0]
        cell.text = (
            f"Submission Date: {investment_firm.submission_date}\n"
            f"{investment_firm.turnaround_requested} turnaround requested\n"
            f"Next Expected Closing: {investment_firm.next_closing}"
        )
        
        # Apply table styles
        apply_table_styles(table)

        debug_info['slide_info'].append({
            'slide': 'Title',
            'status': 'success',
            'time': (dt.now() - slide_start).total_seconds()
        })
    except Exception as e:
        error_msg = f"Error creating title slide: {str(e)}"
        logging.error(error_msg)
        debug_info['errors'].append(error_msg)
        debug_info['slide_info'].append({
            'slide': 'Title',
            'status': 'failed',
            'error': str(e)
        })

    #--- 4 --- Four Box Slide ---
    try:
        slide_start = dt.now()
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['four_block']])
        title = slide.shapes.title
        title.text = "Investment Overview"
        
        # Get the placeholders for each box
        placeholders = {shape.placeholder_format.idx: shape for shape in slide.placeholders}
        
        # Box 1 (Top Left) - Firm Overview
        if 16 in placeholders:  # Typically top-left box
            text_frame = placeholders[16].text_frame
            text_frame.clear()
            p = text_frame.add_paragraph()
            p.text = "Firm Overview"
            p.font.bold = True
            p.font.size = Pt(14)
            
            # Add firm fundamentals
            p = text_frame.add_paragraph()
            p.text = f"Founded: {fundamentals.get('founded_year', 'N/A')}"
            p.font.size = Pt(12)
            
            p = text_frame.add_paragraph()
            p.text = f"HQ: {fundamentals.get('primary_office', 'N/A')}"
            p.font.size = Pt(12)
            
            if fundamentals.get('additional_offices'):
                p = text_frame.add_paragraph()
                p.text = f"Additional Offices: {fundamentals['additional_offices']}"
                p.font.size = Pt(12)
            
            p = text_frame.add_paragraph()
            p.text = f"Ownership: {fundamentals.get('ownership_structure', 'N/A')}"
            p.font.size = Pt(12)
            
            p = text_frame.add_paragraph()
            p.text = f"Employees: {fundamentals.get('total_employees', 'N/A')}"
            p.font.size = Pt(12)
            
            if fundamentals.get('aum'):
                p = text_frame.add_paragraph()
                p.text = f"AUM: ${fundamentals['aum']}"
                p.font.size = Pt(12)

        # Box 2 (Top Right) - Investment Strategy
        if 18 in placeholders:  # Typically top-right box
            text_frame = placeholders[18].text_frame
            text_frame.clear()
            p = text_frame.add_paragraph()
            p.text = "Investment Strategy"
            p.font.bold = True
            p.font.size = Pt(14)
            
            p = text_frame.add_paragraph()
            p.text = fundamentals.get('investment_strategy', 'N/A')
            p.font.size = Pt(12)

        # Box 3 (Bottom Left) - Team Overview
        if 19 in placeholders:  # Typically bottom-left box
            text_frame = placeholders[19].text_frame
            text_frame.clear()
            p = text_frame.add_paragraph()
            p.text = "Team Overview"
            p.font.bold = True
            p.font.size = Pt(14)
            
            # Add team information here
            for member in investment_team[:3]:  # Show first 3 team members
                p = text_frame.add_paragraph()
                p.text = f"{member.get('name', 'N/A')} - {member.get('title', 'N/A')}"
                p.font.size = Pt(12)
            
            if len(investment_team) > 3:
                p = text_frame.add_paragraph()
                p.text = f"+ {len(investment_team) - 3} additional team members"
                p.font.size = Pt(10)
                p.font.italic = True

        # Box 4 (Bottom Right) - Additional Info
        if 20 in placeholders:  # Typically bottom-right box
            text_frame = placeholders[20].text_frame
            text_frame.clear()
            p = text_frame.add_paragraph()
            p.text = "Additional Information"
            p.font.bold = True
            p.font.size = Pt(14)
            
            if fundamentals.get('diversity_status') and fundamentals['diversity_status'].lower() != 'n/a':
                p = text_frame.add_paragraph()
                p.text = f"Diversity Status: {fundamentals['diversity_status']}"
                p.font.size = Pt(12)
            
            if fundamentals.get('website'):
                p = text_frame.add_paragraph()
                p.text = f"Website: {fundamentals['website']}"
                p.font.size = Pt(12)
            
            if fundamentals.get('fund_name'):
                p = text_frame.add_paragraph()
                p.text = f"Current Fund: {fundamentals['fund_name']}"
                p.font.size = Pt(12)

        debug_info['slide_info'].append({
            'slide': 'Investment Overview',
            'status': 'success',
            'time': (dt.now() - slide_start).total_seconds()
        })
    except Exception as e:
        error_msg = f"Error creating four box slide: {str(e)}"
        logging.error(error_msg)
        debug_info['errors'].append(error_msg)
        debug_info['slide_info'].append({
            'slide': 'Investment Overview',
            'status': 'failed',
            'error': str(e)
        })

    #--- 3 --- Investment Team Slide ---
    try:
        slide_start = dt.now()
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['table']])
        title = slide.shapes.title
        title.text = "Investment Team"
        
        # Get the table placeholder (placeholder_idx=13)
        table_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 13:
                table_placeholder = shape
                break
                
        if table_placeholder is None:
            raise ValueError("Could not find table placeholder (idx=13)")
            
        # Create table in placeholder
        rows = len(investment_team) + 1
        cols = 4
        table = table_placeholder.insert_table(rows, cols).table
        
        headers = ["Name", "Title", "Join Year", "Background"]
        for i, header in enumerate(headers):
            table.cell(0, i).text = header
            
        for i, member in enumerate(investment_team, 1):
            table.cell(i, 0).text = member["name"]
            table.cell(i, 1).text = member["title"]
            table.cell(i, 2).text = str(member["join_date"]) if member["join_date"] else "N/A"
            table.cell(i, 3).text = member["background"]
        
        apply_table_styles(table)
        debug_info['slide_info'].append({
            'slide': 'Investment Team',
            'status': 'success',
            'time': (dt.now() - slide_start).total_seconds()
        })
    except Exception as e:
        error_msg = f"Error creating investment team slide: {str(e)}"
        logging.error(error_msg)
        debug_info['errors'].append(error_msg)
        debug_info['slide_info'].append({
            'slide': 'Investment Team',
            'status': 'failed',
            'error': str(e)
        })

    #--- 4 --- Fees and Key Terms Slide ---
    try:
        slide_start = dt.now()
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['table']])
        title = slide.shapes.title
        title.text = "Fees and Key Terms"
        
        terms_data = fees_key_terms[0]
        
        # Get the table placeholder (placeholder_idx=13)
        table_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 13:
                table_placeholder = shape
                break
                
        if table_placeholder is None:
            raise ValueError("Could not find table placeholder (idx=13)")
            
        # Create table in placeholder
        rows = len(terms_data) + 1
        cols = 2
        table = table_placeholder.insert_table(rows, cols).table
        
        headers = ["Term", "Description"]
        for i, header in enumerate(headers):
            table.cell(0, i).text = header
            
        for i, (term, desc) in enumerate(terms_data.items(), 1):
            if term not in ['source', 'source_notes']:
                table.cell(i, 0).text = term.replace('_', ' ').title()
                table.cell(i, 1).text = str(desc)
        
        apply_table_styles(table)
        debug_info['slide_info'].append({
            'slide': 'Fees and Key Terms',
            'status': 'success',
            'time': (dt.now() - slide_start).total_seconds()
        })
    except Exception as e:
        error_msg = f"Error creating fees and key terms slide: {str(e)}"
        logging.error(error_msg)
        debug_info['errors'].append(error_msg)
        debug_info['slide_info'].append({
            'slide': 'Fees and Key Terms',
            'status': 'failed',
            'error': str(e)
        })

    #--- 5 --- Deal History Slide ---
    try:
        slide_start = dt.now()
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['table']])
        title = slide.shapes.title
        title.text = "Deal History"
        
        # Get the table placeholder (placeholder_idx=13)
        table_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 13:
                table_placeholder = shape
                break
                
        if table_placeholder is None:
            raise ValueError("Could not find table placeholder (idx=13)")
            
        # Create table in placeholder with additional columns
        rows = len(deal_history) + 1
        cols = 9  # Increased number of columns
        table = table_placeholder.insert_table(rows, cols).table
        
        # Updated headers to include new fields
        headers = [
            "Firm", 
            "Date", 
            "Type",
            "Cost",
            "FMV",
            "Realized",
            "Proceeds",
            "Multiple",
            "Syndicate Partners"
        ]
        
        for i, header in enumerate(headers):
            table.cell(0, i).text = header
        
        for i, deal in enumerate(deal_history, 1):
            table.cell(i, 0).text = deal["firm"]
            table.cell(i, 1).text = deal["date"]
            table.cell(i, 2).text = deal.get("type_of_investment", "N/A")
            table.cell(i, 3).text = deal.get("cost", "N/A")
            table.cell(i, 4).text = deal.get("fmv", "N/A")
            table.cell(i, 5).text = deal["realized"]
            table.cell(i, 6).text = deal.get("realized_proceeds", "N/A")
            table.cell(i, 7).text = deal.get("multiple", "N/A")
            table.cell(i, 8).text = deal.get("syndicate_partners", "")
        
        apply_table_styles(table)
        debug_info['slide_info'].append({
            'slide': 'Deal History',
            'status': 'success',
            'time': (dt.now() - slide_start).total_seconds()
        })
    except Exception as e:
        error_msg = f"Error creating deal history slide: {str(e)}"
        logging.error(error_msg)
        debug_info['errors'].append(error_msg)
        debug_info['slide_info'].append({
            'slide': 'Deal History',
            'status': 'failed',
            'error': str(e)
        })

    #--- 6 --- Service Providers Slide ---
    try:
        slide_start = dt.now()
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['table']])
        title = slide.shapes.title
        title.text = "Service Providers"
        
        # Get the table placeholder (placeholder_idx=13)
        table_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 13:
                table_placeholder = shape
                break
                
        if table_placeholder is None:
            raise ValueError("Could not find table placeholder (idx=13)")
            
        # Create table in placeholder with additional columns
        rows = len(service_providers) + 1
        cols = 6  # Increased number of columns
        table = table_placeholder.insert_table(rows, cols).table
        
        headers = [
            "Role",
            "Provider",
            "Key Contact",
            "Location",
            "Start Date",
            "Scope of Work"
        ]
        
        for i, header in enumerate(headers):
            table.cell(0, i).text = header
        
        for i, provider in enumerate(service_providers, 1):
            table.cell(i, 0).text = provider["service_type"]
            table.cell(i, 1).text = provider["firm_name"]
            table.cell(i, 2).text = provider.get("key_contact", "N/A")
            table.cell(i, 3).text = provider.get("provider_location", "N/A")
            table.cell(i, 4).text = provider.get("engagement_start_date", "N/A")
            table.cell(i, 5).text = provider.get("scope_of_work", "N/A")
        
        apply_table_styles(table)
        debug_info['slide_info'].append({
            'slide': 'Service Providers',
            'status': 'success',
            'time': (dt.now() - slide_start).total_seconds()
        })
    except Exception as e:
        error_msg = f"Error creating service providers slide: {str(e)}"
        logging.error(error_msg)
        debug_info['errors'].append(error_msg)
        debug_info['slide_info'].append({
            'slide': 'Service Providers',
            'status': 'failed',
            'error': str(e)
        })

    #-- Last -- Debug Info Slide
    if debug:
        try:
            # Create debug info slide
            slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['table_content']])
            title = slide.shapes.title
            title.text = "Debug Info"
            
            # Get the table placeholder (placeholder_idx=13)
            table_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 13:
                    table_placeholder = shape
                    break
                    
            if table_placeholder is None:
                raise ValueError("Could not find table placeholder (idx=13)")
            
            # Create table in placeholder
            rows = len(debug_info['slide_info']) + 2  # +1 for header, +1 for summary
            cols = 3
            table = table_placeholder.insert_table(rows, cols).table
            
            # Set headers
            headers = ["Slide", "Status", "Details"]
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = header
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0, 0, 0)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.font.bold = True
            
            # Add slide info
            for i, info in enumerate(debug_info['slide_info'], start=1):
                table.cell(i, 0).text = info['slide']
                table.cell(i, 1).text = info['status']
                if info['status'] == 'success':
                    table.cell(i, 2).text = f"Created in {info['time']:.2f}s"
                else:
                    table.cell(i, 2).text = f"Error: {info['error']}"
                    for cell in table.rows[i].cells:
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.color.rgb = RGBColor(255, 0, 0)
            
            # Add summary info to placeholder 14
            summary_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 14:
                    summary_placeholder = shape
                    break
                    
            if summary_placeholder is not None:
                summary_frame = summary_placeholder.text_frame
                total_time = (dt.now() - start_time).total_seconds()
                error_count = len(debug_info['errors'])
                
                # Add execution summary
                summary_para = summary_frame.add_paragraph()
                summary_para.text = f"Total execution time: {total_time:.2f}s"
                summary_para.font.size = Pt(10)
                
                error_para = summary_frame.add_paragraph()
                error_para.text = f"Total errors: {error_count}"
                error_para.font.size = Pt(10)
                
                data_para = summary_frame.add_paragraph()
                data_para.text = f"Data loaded: {', '.join(k for k,v in debug_info['data_loaded'].items() if v)}"
                data_para.font.size = Pt(10)

        except Exception as e:
            logging.error(f"Error creating debug info slides: {str(e)}")   

    #------------------------------------------------------------------------------
    # Business Data Slide
    #------------------------------------------------------------------------------
    try:
        data_slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['content']])
        data_slide.shapes.title.text = "Business Object Data"
        
        content = data_slide.placeholders[1].text_frame
        
        # Add business object info
        for attr in ['fundamentals', 'investment_team', 'seed_terms', 'fees_key_terms', 'deal_history', 'service_providers']:
            if hasattr(business, attr):
                attr_para = content.add_paragraph()
                attr_para.text = f"\n{attr.upper()}:"
                attr_para.font.size = Pt(12)
                attr_para.font.bold = True
                
                # Get attribute value and format it
                value = getattr(business, attr)
                if isinstance(value, list):
                    for item in value:
                        item_para = content.add_paragraph()
                        item_para.text = f"- {str(item)}"
                        item_para.font.size = Pt(10)
                else:
                    value_para = content.add_paragraph()
                    value_para.text = str(value)
                    value_para.font.size = Pt(10)
                    
    except Exception as e:
        logging.error(f"Error creating business data slide: {str(e)}")

    #------------------------------------------------------------------------------
    # Template Diagnostics
    #------------------------------------------------------------------------------ 
    try:
        generate_template_diagnostics(prs, layout_index=LAYOUTS['blank'])
    except Exception as e:
        logging.error(f"Error creating template diagnostics slide: {str(e)}")

    # Save to bytes
    try:
        ppt_bytes = BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)
        return ppt_bytes.getvalue()
    except Exception as e:
        logging.error(f"Error saving PowerPoint to bytes: {str(e)}")
        raise

if __name__ == "__main__":
    generate_fund_analysis_ppt()